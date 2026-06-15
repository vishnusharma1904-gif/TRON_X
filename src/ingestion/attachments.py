"""
TRON-X Universal Attachment Processor  (Phase 38)
───────────────────────────────────────────────────
"Support all levels of attachments." A single entry point that turns any
uploaded file into a normalised `Attachment` the chat/agent pipeline can use —
either as injectable text, as multimodal `image_data` blocks (the shape
`orchestrator.chat(image_data=...)` already expects), or both.

Supported kinds (graceful degradation when an optional library is missing):

  documents : .pdf .docx .pptx .xlsx .xls .csv .txt .md .rtf
  data      : .json .yaml .yml .xml .html .htm .log .ini .toml
  code      : .py .js .ts .tsx .jsx .java .c .cpp .h .go .rs .rb .php .sh
              .sql .css .scss .swift .kt .lua .r .m .pl  (and more by ext)
  images    : .png .jpg .jpeg .gif .webp .bmp .tiff .svg   → image_data blocks
  audio     : .mp3 .wav .m4a .ogg .flac .aac                → transcript (if STT)
  video     : .mp4 .mov .mkv .webm .avi                     → metadata (+ optional)
  archives  : .zip                                          → manifest + small text

Design:
  • No hard dependency beyond the stdlib. PyMuPDF / python-docx / openpyxl /
    python-pptx are used *if present*; otherwise the attachment is returned with
    `extracted=False` and a helpful note rather than raising.
  • Text is truncated to `max_chars` (default 24k) so a giant file can't blow
    the context window; `metadata["truncated"]` records when this happens.
  • Pure functions + a thin class. No network, no LLM calls here — extraction
    only. (Transcription is delegated to the existing voice STT module if asked.)

Public API:
    att = process_attachment("/path/to/file.pdf")           # from a path
    att = process_bytes(raw_bytes, filename="report.docx")  # from memory
    att.to_text_block()        # str suitable for prompt injection
    att.image_data             # list[dict] | None for multimodal calls
"""
from __future__ import annotations

import base64
import io
import json
import mimetypes
import os
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Optional

from src.core.logger import log

DEFAULT_MAX_CHARS = 24_000


# ── Extension → kind maps ───────────────────────────────────────────────────────

_IMAGE_EXT = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".tif", ".svg"}
_AUDIO_EXT = {".mp3", ".wav", ".m4a", ".ogg", ".flac", ".aac", ".wma", ".opus"}
_VIDEO_EXT = {".mp4", ".mov", ".mkv", ".webm", ".avi", ".m4v", ".flv", ".wmv"}
_PLAINTEXT_EXT = {".txt", ".md", ".markdown", ".rtf", ".log"}
_DATA_EXT = {".json", ".yaml", ".yml", ".xml", ".html", ".htm", ".ini", ".toml", ".csv", ".tsv"}
_CODE_EXT = {
    ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".c", ".cc", ".cpp", ".cxx",
    ".h", ".hpp", ".go", ".rs", ".rb", ".php", ".sh", ".bash", ".zsh", ".sql",
    ".css", ".scss", ".less", ".swift", ".kt", ".kts", ".lua", ".r", ".m",
    ".mm", ".pl", ".pm", ".dart", ".scala", ".clj", ".ex", ".exs", ".vue",
    ".gradle", ".dockerfile", ".makefile", ".cmake", ".proto", ".graphql",
}
_DOC_EXT = {".pdf", ".docx", ".pptx", ".xlsx", ".xls"}
_ARCHIVE_EXT = {".zip"}


def _kind_for(ext: str, filename: str) -> str:
    ext = ext.lower()
    name = filename.lower()
    if ext in _IMAGE_EXT:
        return "image"
    if ext in _AUDIO_EXT:
        return "audio"
    if ext in _VIDEO_EXT:
        return "video"
    if ext in _DOC_EXT:
        return "document"
    if ext in _ARCHIVE_EXT:
        return "archive"
    if ext in _CODE_EXT or name in ("dockerfile", "makefile"):
        return "code"
    if ext in _DATA_EXT:
        return "data"
    if ext in _PLAINTEXT_EXT:
        return "text"
    return "unknown"


# ── Attachment dataclass ────────────────────────────────────────────────────────

@dataclass
class Attachment:
    """Normalised representation of one processed file."""
    filename: str
    kind: str                                   # image|audio|video|document|code|data|text|archive|unknown
    mime: str = "application/octet-stream"
    size_bytes: int = 0
    text: str = ""                              # extracted text (may be empty)
    image_data: Optional[list[dict]] = None     # multimodal blocks (images only)
    extracted: bool = False                     # did we get usable content?
    note: str = ""                              # human-readable status / hint
    metadata: dict = field(default_factory=dict)

    def to_text_block(self, include_header: bool = True) -> str:
        """Render a prompt-injectable block describing this attachment."""
        header = (f"--- Attachment: {self.filename} "
                  f"({self.kind}, {self.size_bytes} bytes) ---\n")
        if self.kind == "image":
            body = "[image attached — see multimodal content]"
        elif self.text:
            body = self.text
        else:
            body = f"[no extractable text — {self.note or 'unsupported'}]"
        return (header + body) if include_header else body

    def to_dict(self, include_full: bool = False) -> dict:
        d = {
            "filename": self.filename, "kind": self.kind, "mime": self.mime,
            "size_bytes": self.size_bytes, "extracted": self.extracted,
            "note": self.note, "metadata": self.metadata,
            "has_image_data": bool(self.image_data),
            "text_preview": self.text[:500],
            "text_len": len(self.text),
        }
        if include_full:
            d["text"] = self.text
            d["image_data"] = self.image_data
        return d


# ── Per-kind extractors (all pure, all degrade gracefully) ──────────────────────

def _truncate(text: str, max_chars: int) -> tuple[str, bool]:
    if len(text) <= max_chars:
        return text, False
    return text[:max_chars] + f"\n…[truncated, {len(text) - max_chars} more chars]", True


def _decode_text(raw: bytes) -> str:
    """Best-effort decode with a couple of fallbacks."""
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            return raw.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return raw.decode("utf-8", errors="replace")


def _extract_pdf(raw: bytes) -> tuple[str, dict]:
    try:
        import fitz  # PyMuPDF
    except Exception:
        return "", {"error": "PyMuPDF not installed (pip install pymupdf)"}
    try:
        doc = fitz.open(stream=raw, filetype="pdf")
        pages = [p.get_text() for p in doc]
        meta = {"pages": len(pages)}
        return "\n\n".join(pages).strip(), meta
    except Exception as e:
        return "", {"error": f"pdf parse failed: {e}"}


def _extract_docx(raw: bytes) -> tuple[str, dict]:
    try:
        from docx import Document  # python-docx
    except Exception:
        return "", {"error": "python-docx not installed (pip install python-docx)"}
    try:
        doc = Document(io.BytesIO(raw))
        paras = [p.text for p in doc.paragraphs if p.text.strip()]
        # include table cells too
        for tbl in getattr(doc, "tables", []):
            for row in tbl.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    paras.append(" | ".join(cells))
        return "\n".join(paras).strip(), {"paragraphs": len(paras)}
    except Exception as e:
        return "", {"error": f"docx parse failed: {e}"}


def _extract_pptx(raw: bytes) -> tuple[str, dict]:
    try:
        from pptx import Presentation  # python-pptx
    except Exception:
        return "", {"error": "python-pptx not installed (pip install python-pptx)"}
    try:
        prs = Presentation(io.BytesIO(raw))
        out = []
        for i, slide in enumerate(prs.slides, 1):
            out.append(f"[Slide {i}]")
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs)
                        if line.strip():
                            out.append(line)
        return "\n".join(out).strip(), {"slides": len(prs.slides)}
    except Exception as e:
        return "", {"error": f"pptx parse failed: {e}"}


def _extract_xlsx(raw: bytes) -> tuple[str, dict]:
    try:
        import openpyxl
    except Exception:
        return "", {"error": "openpyxl not installed (pip install openpyxl)"}
    try:
        wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        out = []
        for ws in wb.worksheets:
            out.append(f"[Sheet: {ws.title}]")
            for row in ws.iter_rows(values_only=True):
                cells = ["" if c is None else str(c) for c in row]
                if any(cells):
                    out.append(" | ".join(cells))
        return "\n".join(out).strip(), {"sheets": len(wb.worksheets)}
    except Exception as e:
        return "", {"error": f"xlsx parse failed: {e}"}


def _extract_csv(raw: bytes) -> tuple[str, dict]:
    import csv
    text = _decode_text(raw)
    try:
        rows = list(csv.reader(io.StringIO(text)))
        out = [" | ".join(r) for r in rows]
        return "\n".join(out).strip(), {"rows": len(rows)}
    except Exception:
        return text.strip(), {"rows": text.count("\n") + 1}


def _extract_json(raw: bytes) -> tuple[str, dict]:
    text = _decode_text(raw)
    try:
        obj = json.loads(text)
        pretty = json.dumps(obj, indent=2, ensure_ascii=False)
        keys = list(obj.keys()) if isinstance(obj, dict) else None
        return pretty.strip(), {"valid_json": True, "top_keys": keys}
    except Exception:
        return text.strip(), {"valid_json": False}


def _extract_archive(raw: bytes, max_chars: int) -> tuple[str, dict]:
    import zipfile
    try:
        zf = zipfile.ZipFile(io.BytesIO(raw))
    except Exception as e:
        return "", {"error": f"not a valid zip: {e}"}
    names = zf.namelist()
    manifest = [f"Archive contains {len(names)} entries:"]
    manifest += [f"  - {n}" for n in names[:200]]
    # Inline small text-ish files
    budget = max_chars
    previews = []
    for n in names:
        ext = Path(n).suffix.lower()
        if ext in _PLAINTEXT_EXT | _DATA_EXT | _CODE_EXT and budget > 0:
            try:
                inner = zf.read(n)
                if len(inner) > 50_000:
                    continue
                snippet = _decode_text(inner)[:2000]
                block = f"\n----- {n} -----\n{snippet}"
                previews.append(block)
                budget -= len(block)
            except Exception:
                continue
    return "\n".join(manifest + previews).strip(), {"entries": len(names)}


# ── Processor ───────────────────────────────────────────────────────────────────

class AttachmentProcessor:
    """Routes a file (by bytes + filename) to the right extractor."""

    def __init__(self, max_chars: int = DEFAULT_MAX_CHARS, enable_stt: bool = False):
        self.max_chars = int(max_chars)
        self.enable_stt = bool(enable_stt)

    # -- core -------------------------------------------------------------------

    def process_bytes(self, raw: bytes, filename: str) -> Attachment:
        ext = Path(filename).suffix.lower()
        kind = _kind_for(ext, filename)
        mime = mimetypes.guess_type(filename)[0] or "application/octet-stream"
        att = Attachment(filename=filename, kind=kind, mime=mime, size_bytes=len(raw))

        try:
            if kind == "image":
                self._handle_image(att, raw, mime)
            elif kind == "document":
                self._handle_document(att, raw, ext)
            elif kind == "archive":
                text, meta = _extract_archive(raw, self.max_chars)
                self._set_text(att, text, meta)
            elif kind in ("code", "text"):
                self._set_text(att, _decode_text(raw), {})
            elif kind == "data":
                self._handle_data(att, raw, ext)
            elif kind == "audio":
                self._handle_audio(att, raw)
            elif kind == "video":
                att.note = ("video accepted; transcript/frame extraction not run in this "
                            "context. Provide audio track or keyframes for analysis.")
                att.metadata["needs_external"] = True
            else:  # unknown — try a text decode, else mark binary
                decoded = _decode_text(raw)
                printable = sum(c.isprintable() or c in "\r\n\t" for c in decoded[:2000])
                if decoded and printable / max(1, len(decoded[:2000])) > 0.85:
                    att.kind = "text"
                    self._set_text(att, decoded, {"sniffed": True})
                else:
                    att.note = "binary file of unrecognised type; no text extracted."
        except Exception as e:
            att.note = f"processing error: {e}"
            log.warning(f"[attachments] {filename}: {e}")

        return att

    def process_path(self, path: str) -> Attachment:
        p = Path(path)
        if not p.exists():
            return Attachment(filename=p.name, kind="unknown",
                              note=f"file not found: {path}")
        raw = p.read_bytes()
        return self.process_bytes(raw, p.name)

    # -- handlers ---------------------------------------------------------------

    def _set_text(self, att: Attachment, text: str, meta: dict) -> None:
        if meta.get("error") and not text:
            att.extracted = False
            att.note = meta["error"]
            att.metadata.update(meta)
            return
        truncated_text, was_trunc = _truncate(text, self.max_chars)
        att.text = truncated_text
        att.extracted = bool(text.strip())
        att.metadata.update(meta)
        att.metadata["truncated"] = was_trunc
        if not att.extracted and not att.note:
            att.note = "no text content found."

    def _handle_image(self, att: Attachment, raw: bytes, mime: str) -> None:
        if att.filename.lower().endswith(".svg"):
            # SVG is XML/text — give the model the source AND mark as image.
            att.text = _truncate(_decode_text(raw), self.max_chars)[0]
            att.extracted = True
            att.note = "SVG provided as source text."
            return
        b64 = base64.b64encode(raw).decode()
        if mime == "application/octet-stream":
            mime = "image/png"
        att.image_data = [{"type": "image_url",
                           "image_url": {"url": f"data:{mime};base64,{b64}"}}]
        att.extracted = True
        att.note = "image ready for multimodal model."

    def _handle_document(self, att: Attachment, raw: bytes, ext: str) -> None:
        if ext == ".pdf":
            text, meta = _extract_pdf(raw)
        elif ext == ".docx":
            text, meta = _extract_docx(raw)
        elif ext == ".pptx":
            text, meta = _extract_pptx(raw)
        elif ext in (".xlsx", ".xls"):
            text, meta = _extract_xlsx(raw)
        else:
            text, meta = "", {"error": f"unsupported document type {ext}"}
        self._set_text(att, text, meta)

    def _handle_data(self, att: Attachment, raw: bytes, ext: str) -> None:
        if ext == ".json":
            text, meta = _extract_json(raw)
        elif ext in (".csv", ".tsv"):
            text, meta = _extract_csv(raw)
        else:  # yaml/xml/html/ini/toml — keep as text
            text, meta = _decode_text(raw), {}
        self._set_text(att, text, meta)

    def _handle_audio(self, att: Attachment, raw: bytes) -> None:
        if not self.enable_stt:
            att.note = ("audio accepted; speech-to-text not enabled for this call "
                        "(set enable_stt=True to transcribe via the voice module).")
            att.metadata["needs_stt"] = True
            return
        try:
            from src.voice.stt import transcribe_bytes  # optional, may not exist
            transcript = transcribe_bytes(raw)           # pragma: no cover
            self._set_text(att, transcript or "", {"source": "stt"})
            att.note = "audio transcribed via voice module."
        except Exception as e:
            att.note = f"audio transcription unavailable: {e}"
            att.metadata["needs_stt"] = True


# ── Module-level convenience ────────────────────────────────────────────────────

def process_bytes(raw: bytes, filename: str, max_chars: int = DEFAULT_MAX_CHARS,
                  enable_stt: bool = False) -> Attachment:
    return AttachmentProcessor(max_chars=max_chars, enable_stt=enable_stt).process_bytes(raw, filename)


def process_attachment(path: str, max_chars: int = DEFAULT_MAX_CHARS,
                       enable_stt: bool = False) -> Attachment:
    return AttachmentProcessor(max_chars=max_chars, enable_stt=enable_stt).process_path(path)


def process_many(files: list[tuple[bytes, str]], max_chars: int = DEFAULT_MAX_CHARS,
                 enable_stt: bool = False) -> list[Attachment]:
    """Process a batch of (bytes, filename) pairs."""
    proc = AttachmentProcessor(max_chars=max_chars, enable_stt=enable_stt)
    return [proc.process_bytes(raw, name) for raw, name in files]


def merge_for_prompt(attachments: list[Attachment]) -> tuple[str, Optional[list[dict]]]:
    """
    Combine processed attachments into (text_block, image_data) ready to pass to
    orchestrator.chat(extra_system=..., image_data=...). Image blocks from all
    image attachments are concatenated; text blocks are joined with separators.
    """
    text_parts = [a.to_text_block() for a in attachments if a.kind != "image" or a.text]
    images: list[dict] = []
    for a in attachments:
        if a.image_data:
            images.extend(a.image_data)
    text_block = "\n\n".join(tp for tp in text_parts if tp.strip())
    return text_block, (images or None)
