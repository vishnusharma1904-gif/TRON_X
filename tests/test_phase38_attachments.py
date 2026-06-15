"""
Phase 38 — Universal Attachment Processor
Standalone test (conventions per TRONX_PHASE21-36_PROGRESS_HANDOFF.md §3).
Builds real PDF/DOCX/PPTX/XLSX/ZIP fixtures in-memory where libraries are
available; gracefully skips a format if its library is missing (mirroring the
processor's own graceful degradation).
"""
from __future__ import annotations

import base64
import io
import json
import os
import sys
import zipfile
from unittest.mock import MagicMock

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)
os.chdir(ROOT)

if "chromadb" not in sys.modules:
    sys.modules["chromadb"] = MagicMock()
    sys.modules["chromadb.config"] = MagicMock()

from src.ingestion.attachments import (  # noqa: E402
    AttachmentProcessor,
    merge_for_prompt,
    process_bytes,
    process_many,
)

PASS = 0
FAIL = 0
SKIP = 0


def check(name: str, cond: bool, detail: str = "") -> None:
    global PASS, FAIL
    if cond:
        PASS += 1
        print(f"  PASS  {name}")
    else:
        FAIL += 1
        print(f"  FAIL  {name}  {detail}")


def skip(name: str, why: str) -> None:
    global SKIP
    SKIP += 1
    print(f"  SKIP  {name}  ({why})")


# 1x1 transparent PNG
_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR4nGNgYGBg"
    "AAAABQABh6FO1AAAAABJRU5ErkJggg=="
)


# =============================================================================
print("\n=== kind detection ===")
# =============================================================================

cases = {
    "report.pdf": "document", "deck.pptx": "document", "sheet.xlsx": "document",
    "notes.txt": "text", "readme.md": "text", "photo.jpg": "image",
    "song.mp3": "audio", "clip.mp4": "video", "data.json": "data",
    "table.csv": "data", "app.py": "code", "archive.zip": "archive",
    "mystery.xyz": "unknown",
}
ok = all(process_bytes(b"", n).kind == k or (k == "unknown")
         for n, k in cases.items())
for name_, kind_ in cases.items():
    att = process_bytes(b"x", name_)
    if kind_ == "unknown":
        continue  # unknown may be sniffed to text — checked separately below
    check(f"{name_} -> {kind_}", att.kind == kind_, att.kind)


# =============================================================================
print("\n=== text / code / data extraction ===")
# =============================================================================

att = process_bytes("hello world\nsecond line".encode(), "notes.txt")
check("txt extracted", att.extracted and "second line" in att.text, att.to_dict())

att = process_bytes(b"def f():\n    return 42\n", "app.py")
check("code extracted verbatim", "return 42" in att.text, att.text)

att = process_bytes(json.dumps({"a": 1, "b": [2, 3]}).encode(), "data.json")
check("json pretty-printed", att.extracted and '"a": 1' in att.text, att.text)
check("json metadata top_keys", att.metadata.get("top_keys") == ["a", "b"], att.metadata)

att = process_bytes(b"col1,col2\nx,y\n", "table.csv")
check("csv rows extracted", "col1 | col2" in att.text and att.metadata.get("rows") == 2,
      att.to_dict())

att = process_bytes("ünïcode tëxt".encode("utf-8"), "u.txt")
check("utf-8 decode", "ünïcode" in att.text, att.text)

att = process_bytes("latin1 café".encode("latin-1"), "l.txt")
check("latin-1 fallback decode", "caf" in att.text, att.text)


# =============================================================================
print("\n=== truncation ===")
# =============================================================================

big = ("word " * 2000).encode()
att = AttachmentProcessor(max_chars=500).process_bytes(big, "big.txt")
check("text truncated to budget", len(att.text) < 600, len(att.text))
check("truncated flag set", att.metadata.get("truncated") is True, att.metadata)


# =============================================================================
print("\n=== images ===")
# =============================================================================

att = process_bytes(_PNG, "photo.png")
check("image_data block built", bool(att.image_data), att.to_dict())
check("data-url shape", att.image_data[0]["image_url"]["url"].startswith("data:image/png;base64,"),
      att.image_data[0]["image_url"]["url"][:40])
check("image marked extracted", att.extracted is True)

att = process_bytes(b"<svg xmlns='http://www.w3.org/2000/svg'><rect/></svg>", "icon.svg")
check("svg returned as source text", att.extracted and "<rect/>" in att.text, att.to_dict())


# =============================================================================
print("\n=== documents (real fixtures, skip if lib missing) ===")
# =============================================================================

# PDF via PyMuPDF
try:
    import fitz
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "TRONX PDF FIXTURE")
    pdf_bytes = doc.tobytes()
    att = process_bytes(pdf_bytes, "fixture.pdf")
    check("pdf text extracted", att.extracted and "TRONX PDF FIXTURE" in att.text, att.to_dict())
    check("pdf page count metadata", att.metadata.get("pages") == 1, att.metadata)
except ImportError:
    skip("pdf extraction", "PyMuPDF not installed")

# DOCX via python-docx
try:
    from docx import Document
    buf = io.BytesIO()
    d = Document()
    d.add_paragraph("TRONX DOCX FIXTURE")
    tbl = d.add_table(rows=1, cols=2)
    tbl.rows[0].cells[0].text = "cellA"
    tbl.rows[0].cells[1].text = "cellB"
    d.save(buf)
    att = process_bytes(buf.getvalue(), "fixture.docx")
    check("docx paragraphs extracted", "TRONX DOCX FIXTURE" in att.text, att.to_dict())
    check("docx table cells extracted", "cellA | cellB" in att.text, att.text)
except ImportError:
    skip("docx extraction", "python-docx not installed")

# PPTX via python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches
    buf = io.BytesIO()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tb.text_frame.text = "TRONX PPTX FIXTURE"
    prs.save(buf)
    att = process_bytes(buf.getvalue(), "fixture.pptx")
    check("pptx text extracted", "TRONX PPTX FIXTURE" in att.text, att.to_dict())
    check("pptx slide marker", "[Slide 1]" in att.text, att.text)
except ImportError:
    skip("pptx extraction", "python-pptx not installed")

# XLSX via openpyxl
try:
    import openpyxl
    buf = io.BytesIO()
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["name", "value"])
    ws.append(["tronx", 42])
    wb.save(buf)
    att = process_bytes(buf.getvalue(), "fixture.xlsx")
    check("xlsx cells extracted", "tronx | 42" in att.text, att.to_dict())
    check("xlsx sheet marker", "[Sheet: Data]" in att.text, att.text)
except ImportError:
    skip("xlsx extraction", "openpyxl not installed")


# =============================================================================
print("\n=== archives ===")
# =============================================================================

buf = io.BytesIO()
with zipfile.ZipFile(buf, "w") as zf:
    zf.writestr("readme.md", "# Inner readme")
    zf.writestr("code/util.py", "def util(): pass")
    zf.writestr("blob.bin", b"\x00\x01\x02")
att = process_bytes(buf.getvalue(), "bundle.zip")
check("zip manifest listed", "3 entries" in att.text, att.text[:200])
check("inner text inlined", "# Inner readme" in att.text and "def util" in att.text,
      att.text[:400])
check("zip metadata entries", att.metadata.get("entries") == 3, att.metadata)

att = process_bytes(b"not a zip at all", "broken.zip")
check("invalid zip degrades", att.extracted is False and "zip" in (att.note + str(att.metadata)),
      att.to_dict())


# =============================================================================
print("\n=== audio / video / unknown ===")
# =============================================================================

att = process_bytes(b"\xff\xfb\x90fake-mp3", "song.mp3")
check("audio accepted, needs_stt flagged", att.metadata.get("needs_stt") is True, att.to_dict())
check("audio not extracted by default", att.extracted is False)

att = process_bytes(b"\x00fakevideo", "clip.mp4")
check("video accepted with note", att.kind == "video" and att.note, att.to_dict())

att = process_bytes(b"plain readable content here", "mystery.xyz")
check("unknown text sniffed", att.kind == "text" and att.extracted, att.to_dict())

att = process_bytes(bytes(range(256)) * 4, "mystery.bin")
check("unknown binary not extracted", att.extracted is False, att.to_dict())


# =============================================================================
print("\n=== prompt assembly ===")
# =============================================================================

atts = process_many([
    (b"hello doc", "a.txt"),
    (_PNG, "b.png"),
])
text_block, image_data = merge_for_prompt(atts)
check("text block contains file header", "Attachment: a.txt" in text_block, text_block[:120])
check("image_data propagated", image_data and len(image_data) == 1, image_data)

block = atts[0].to_text_block()
check("to_text_block has header+body", block.startswith("--- Attachment: a.txt")
      and "hello doc" in block, block)

# error path: missing file
proc = AttachmentProcessor()
att = proc.process_path("/definitely/not/here.pdf")
check("missing path -> note, no raise", "not found" in att.note, att.note)


# =============================================================================
print(f"\n{'=' * 60}\nRESULT: {PASS} passed, {FAIL} failed, {SKIP} skipped\n{'=' * 60}")
sys.exit(1 if FAIL else 0)
